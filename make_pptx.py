from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── color palette ────────────────────────────────────────────────────────────
BLUE_DARK  = RGBColor(0x1E, 0x3A, 0x5F)   # sidebar navy
BLUE_MID   = RGBColor(0x1D, 0x4E, 0xD8)   # btn-primary
BLUE_LIGHT = RGBColor(0xDB, 0xEA, 0xFE)   # badge bg
TEAL       = RGBColor(0x0D, 0x94, 0x88)
GREEN      = RGBColor(0x16, 0xA3, 0x4A)
ORANGE     = RGBColor(0xEA, 0x58, 0x0C)
RED        = RGBColor(0xDC, 0x26, 0x26)
GRAY_DARK  = RGBColor(0x1F, 0x29, 0x37)
GRAY_MID   = RGBColor(0x6B, 0x72, 0x80)
GRAY_LIGHT = RGBColor(0xF3, 0xF4, 0xF6)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
YELLOW     = RGBColor(0xFA, 0xCC, 0x15)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ─── helpers ──────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None, radius=False):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, x, y, w, h,
             size=14, bold=False, color=GRAY_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return txb


def add_para(tf, text, size=13, bold=False, color=GRAY_DARK,
             align=PP_ALIGN.LEFT, space_before=0, italic=False):
    p = tf.add_paragraph()
    p.alignment = align
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.italic = italic
    return p


def sidebar(slide):
    """Draw left sidebar identical to app sidebar."""
    add_rect(slide, 0, 0, 2.2, 7.5, fill=BLUE_DARK)
    # logo area
    add_rect(slide, 0, 0, 2.2, 0.85, fill=RGBColor(0x17, 0x30, 0x52))
    add_text(slide, "❄  Air System", 0.15, 0.12, 2.0, 0.38,
             size=13, bold=True, color=WHITE)
    add_text(slide, "Technical Water", 0.15, 0.46, 2.0, 0.28,
             size=9, color=RGBColor(0x93, 0xC5, 0xFD))

    nav = [
        ("Dashboard",    False),
        ("แจ้งซ่อม",      False),
        ("PM Schedule",  False),
        ("ใบงาน",         False),
        ("Master Data",  False),
        ("ผู้ใช้งาน",      False),
    ]
    for i, (label, active) in enumerate(nav):
        yy = 0.95 + i * 0.53
        if active:
            add_rect(slide, 0.12, yy - 0.05, 1.96, 0.44,
                     fill=RGBColor(0x1D, 0x4E, 0xD8))
        add_text(slide, label, 0.25, yy, 1.8, 0.38,
                 size=11, color=WHITE if active else RGBColor(0xBF, 0xDB, 0xFE))

    add_text(slide, "admin  •  Admin", 0.15, 6.8, 2.0, 0.35,
             size=9, color=RGBColor(0x93, 0xC5, 0xFD))


def header_bar(slide, title, subtitle=""):
    add_rect(slide, 2.2, 0, 11.13, 0.65, fill=WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.5)
    add_text(slide, title, 2.45, 0.1, 8.0, 0.45,
             size=15, bold=True, color=GRAY_DARK)
    if subtitle:
        add_text(slide, subtitle, 2.45, 0.38, 8.0, 0.25,
                 size=9, color=GRAY_MID)


def slide_title_only(title, subtitle=""):
    """Title slide — full navy bg."""
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
    # decorative strip
    add_rect(sl, 0, 6.3, 13.33, 0.05, fill=BLUE_MID)
    add_rect(sl, 0, 6.35, 13.33, 1.15, fill=RGBColor(0x17, 0x30, 0x52))
    add_text(sl, "❄", 1.0, 1.6, 1.0, 1.0,
             size=48, color=RGBColor(0x60, 0xA5, 0xFA))
    add_text(sl, title, 1.0, 2.55, 11.0, 1.2,
             size=36, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(sl, subtitle, 1.0, 3.75, 11.0, 0.6,
                 size=18, color=RGBColor(0x93, 0xC5, 0xFD),
                 align=PP_ALIGN.LEFT)
    add_text(sl, "บริษัท Technical Water Co.,Ltd", 1.0, 6.5, 8.0, 0.4,
             size=11, color=RGBColor(0x93, 0xC5, 0xFD))
    return sl


def section_divider(title, color=BLUE_MID):
    sl = prs.slides.add_slide(BLANK)
    add_rect(sl, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
    add_rect(sl, 0, 3.2, 0.18, 1.1, fill=color)
    add_text(sl, title, 0.45, 3.2, 12.0, 1.1,
             size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    return sl


def new_slide(title="", subtitle=""):
    sl = prs.slides.add_slide(BLANK)
    # white bg
    add_rect(sl, 0, 0, 13.33, 7.5, fill=WHITE)
    sidebar(sl)
    header_bar(sl, title, subtitle)
    return sl


def stat_card(slide, x, y, w, h, label, value, color, icon=""):
    add_rect(slide, x, y, w, h, fill=WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.8)
    # color accent top
    add_rect(slide, x, y, w, 0.07, fill=color)
    if icon:
        add_text(slide, icon, x + 0.12, y + 0.15, 0.45, 0.45, size=20, color=color)
    add_text(slide, value, x + 0.62, y + 0.12, w - 0.75, 0.5,
             size=22, bold=True, color=color, align=PP_ALIGN.LEFT)
    add_text(slide, label, x + 0.12, y + 0.65, w - 0.25, 0.35,
             size=10, color=GRAY_MID)


def role_badge(slide, x, y, label, bg, fg=WHITE):
    add_rect(slide, x, y, 1.35, 0.32, fill=bg)
    add_text(slide, label, x + 0.08, y + 0.05, 1.2, 0.24,
             size=10, bold=True, color=fg, align=PP_ALIGN.CENTER)


def bullet_box(slide, x, y, w, h, title, bullets, title_color=BLUE_MID,
               bg=GRAY_LIGHT, icon=""):
    add_rect(slide, x, y, w, h, fill=bg,
             line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.5)
    ty = y + 0.12
    add_text(slide, (icon + "  " if icon else "") + title,
             x + 0.18, ty, w - 0.3, 0.38,
             size=12, bold=True, color=title_color)
    for i, b in enumerate(bullets):
        add_text(slide, "  •  " + b, x + 0.18, ty + 0.38 + i * 0.34,
                 w - 0.3, 0.32, size=10, color=GRAY_DARK)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Cover
# ══════════════════════════════════════════════════════════════════════════════
slide_title_only(
    "Air System",
    "คู่มือการใช้งานระบบบริหารจัดการแอร์  |  Air Conditioning Maintenance System"
)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — ภาพรวมระบบ
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("ภาพรวมระบบ", "Air System ครอบคลุม 6 โมดูลหลัก")

modules = [
    ("🏠", "Dashboard",   "สรุปสถิติ + กราฟ",                BLUE_MID),
    ("🔧", "แจ้งซ่อม",    "แจ้งปัญหา → ติดตามสถานะ",          ORANGE),
    ("📅", "PM Schedule", "แผนซ่อมบำรุงตามรอบ",              TEAL),
    ("📋", "ใบงาน",       "สร้าง / บันทึกงาน / อนุมัติ",     BLUE_MID),
    ("🗄", "Master Data",  "จัดการโรงพยาบาล/อาคาร/เครื่องแอร์", GRAY_MID),
    ("👤", "ผู้ใช้งาน",   "จัดการ account + สิทธิ์",         RED),
]

for i, (icon, name, desc, col) in enumerate(modules):
    col_i = i % 3
    row_i = i // 3
    x = 2.35 + col_i * 3.6
    y = 0.85 + row_i * 2.9
    add_rect(sl, x, y, 3.3, 2.55, fill=WHITE,
             line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.8)
    add_rect(sl, x, y, 3.3, 0.08, fill=col)
    add_text(sl, icon, x + 0.2, y + 0.18, 0.6, 0.55, size=24, color=col)
    add_text(sl, name, x + 0.85, y + 0.2, 2.2, 0.42,
             size=14, bold=True, color=GRAY_DARK)
    add_text(sl, desc, x + 0.2, y + 0.72, 2.9, 0.45,
             size=10, color=GRAY_MID)

# Role boxes
add_text(sl, "สิทธิ์การใช้งาน", 2.35, 6.75, 2.5, 0.35,
         size=10, bold=True, color=GRAY_MID)
roles = [
    ("Admin", BLUE_DARK),
    ("Owner", TEAL),
    ("Tech",  GREEN),
]
for i, (r, c) in enumerate(roles):
    role_badge(sl, 4.4 + i * 1.6, 6.72, r, c)
add_text(sl, "Admin = ทุกสิทธิ์   |   Owner = ดู+อนุมัติ+Master   |   Tech = สร้าง+บันทึกงาน",
         2.35, 7.1, 10.5, 0.3, size=9, color=GRAY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Login
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("การเข้าสู่ระบบ (Login)")

# mock login card
add_rect(sl, 4.5, 1.0, 4.3, 5.3, fill=WHITE,
         line=RGBColor(0xE5, 0xE7, 0xEB), line_w=1)
add_rect(sl, 4.5, 1.0, 4.3, 0.08, fill=BLUE_MID)

add_text(sl, "❄  Air System", 4.7, 1.2, 3.9, 0.5,
         size=18, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)
add_text(sl, "เข้าสู่ระบบ", 4.7, 1.75, 3.9, 0.38,
         size=13, color=GRAY_MID, align=PP_ALIGN.CENTER)

# username field
add_rect(sl, 4.75, 2.35, 3.8, 0.52, fill=WHITE,
         line=RGBColor(0xD1, 0xD5, 0xDB), line_w=1)
add_text(sl, "Username", 4.95, 2.44, 3.5, 0.32, size=10, color=GRAY_MID)

# password field
add_rect(sl, 4.75, 3.05, 3.8, 0.52, fill=WHITE,
         line=RGBColor(0xD1, 0xD5, 0xDB), line_w=1)
add_text(sl, "Password", 4.95, 3.14, 3.5, 0.32, size=10, color=GRAY_MID)

# login btn
add_rect(sl, 4.75, 3.75, 3.8, 0.52, fill=BLUE_MID)
add_text(sl, "เข้าสู่ระบบ", 4.75, 3.82, 3.8, 0.38,
         size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# note
add_text(sl, "* Session จะถูกบันทึกใน Browser\n  ไม่ต้อง login ซ้ำเมื่อเปิดหน้าใหม่",
         4.75, 4.45, 3.8, 0.6, size=9, color=GRAY_MID, align=PP_ALIGN.CENTER)

# side notes
bullet_box(sl, 2.35, 1.3, 1.85, 2.5, "บัญชีผู้ใช้", [
    "admin — Admin",
    "owner — Owner",
    "tech1 — Tech",
    "tech2 — Tech",
], bg=WHITE, icon="👤")

bullet_box(sl, 9.1, 1.3, 4.0, 3.5, "หลังเข้าสู่ระบบ", [
    "ระบบ redirect → Dashboard",
    "Sidebar แสดงเมนูตาม Role",
    "Token เก็บใน localStorage",
    "401 → logout อัตโนมัติ",
    "ออกจากระบบ: ปุ่มล่างซ้าย Sidebar",
], bg=GRAY_LIGHT, icon="ℹ️")

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Dashboard
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("Dashboard", "สรุปสถิติภาพรวมระบบ — แสดงอัตโนมัติหลังเข้าสู่ระบบ")

# stat cards row 1
cards1 = [
    ("กำลังดำเนินการ", "12", BLUE_MID,  "📋"),
    ("รออนุมัติ",       "3",  YELLOW,    "⏳"),
    ("อนุมัติแล้ว",     "87", GREEN,     "✅"),
    ("แจ้งซ่อมค้าง",   "5",  RED,       "🔧"),
    ("โรงพยาบาล",      "4",  TEAL,      "🏥"),
    ("เครื่องแอร์",    "248","#6366F1", "❄"),
]
for i, (lbl, val, col, icon) in enumerate(cards1):
    x = 2.35 + i * 1.82
    try:
        c = RGBColor.from_string(col.replace("#","")) if col.startswith("#") else col
    except:
        c = RGBColor(0x63, 0x66, 0xF1)
    stat_card(sl, x, 0.75, 1.72, 1.18, lbl, val, c, icon)

# chart mock
add_rect(sl, 2.35, 2.12, 5.5, 3.0, fill=WHITE,
         line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.8)
add_text(sl, "ใบงานรายเดือน (6 เดือนล่าสุด)", 2.55, 2.22, 5.1, 0.38,
         size=11, bold=True, color=GRAY_DARK)
months = ["ธ.ค.", "ม.ค.", "ก.พ.", "มี.ค.", "เม.ย.", "พ.ค."]
heights = [0.8, 1.1, 0.6, 1.4, 0.9, 1.2]
bar_w = 0.5
for i, (m, bh) in enumerate(zip(months, heights)):
    bx = 2.6 + i * 0.83
    by = 2.6 + (1.4 - bh)
    add_rect(sl, bx, by, bar_w, bh, fill=BLUE_MID)
    add_text(sl, m, bx, by + bh + 0.04, bar_w, 0.22,
             size=8, color=GRAY_MID, align=PP_ALIGN.CENTER)

# type breakdown mock
add_rect(sl, 8.05, 2.12, 5.1, 3.0, fill=WHITE,
         line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.8)
add_text(sl, "ใบงานตามประเภท", 8.25, 2.22, 4.7, 0.38,
         size=11, bold=True, color=GRAY_DARK)
types = [
    ("ล้างใหญ่ (Major)", 65, BLUE_MID),
    ("ล้างย่อย (Minor)", 28, TEAL),
    ("ล้างคอยล์พัดลม (Fan)", 7, GREEN),
]
for i, (t, pct, c) in enumerate(types):
    y = 2.75 + i * 0.72
    add_rect(sl, 8.25, y, 3.5, 0.3, fill=GRAY_LIGHT)
    add_rect(sl, 8.25, y, 3.5 * pct / 100, 0.3, fill=c)
    add_text(sl, t, 8.25, y - 0.24, 2.8, 0.24, size=9, color=GRAY_DARK)
    add_text(sl, f"{pct}%", 11.65, y - 0.24, 0.5, 0.24,
             size=9, bold=True, color=c)

# recent WO table
add_rect(sl, 2.35, 5.3, 10.8, 1.95, fill=WHITE,
         line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.8)
add_text(sl, "ใบงานล่าสุด", 2.55, 5.38, 4.0, 0.35,
         size=11, bold=True, color=GRAY_DARK)
# table header
add_rect(sl, 2.35, 5.75, 10.8, 0.35, fill=GRAY_LIGHT)
for txt, xx, ww in [("เลขที่ใบงาน", 2.45, 1.6), ("โรงพยาบาล", 4.1, 2.0),
                     ("ประเภท", 6.15, 1.5), ("สถานะ", 7.7, 1.5),
                     ("ช่างเทคนิค", 9.25, 1.8), ("วันที่", 11.1, 0.9)]:
    add_text(sl, txt, xx, 5.79, ww, 0.28, size=8, bold=True, color=GRAY_MID)
rows = [
    ("WO-001", "รพ.กรุงเทพ", "Major", "อนุมัติ", "tech1", "20/05"),
    ("WO-002", "รพ.สมิติเวช", "Minor", "กำลังทำ", "tech2", "22/05"),
    ("WO-003", "รพ.วิภาวดี", "Fan",   "รออนุมัติ","tech1", "24/05"),
]
for ri, row in enumerate(rows):
    yy = 6.12 + ri * 0.32
    for ci, (txt, xx, ww) in enumerate(zip(row,
            [2.45, 4.1, 6.15, 7.7, 9.25, 11.1],
            [1.6,  2.0, 1.5,  1.5,  1.8,  0.9])):
        add_text(sl, txt, xx, yy, ww, 0.28, size=9, color=GRAY_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — แจ้งซ่อม (Repair Logs)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("แจ้งซ่อม (Repair Logs)", "บันทึกการแจ้งซ่อม + ติดตามสถานะ")

# left: flow
add_rect(sl, 2.35, 0.75, 3.6, 5.9, fill=GRAY_LIGHT,
         line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.5)
add_text(sl, "ขั้นตอนการแจ้งซ่อม", 2.5, 0.82, 3.3, 0.38,
         size=11, bold=True, color=BLUE_DARK)

steps = [
    (BLUE_MID,  "1", "กด + แจ้งซ่อมใหม่"),
    (TEAL,      "2", "เลือกโรงพยาบาล"),
    (TEAL,      "3", "ค้นหาและเลือกเครื่องแอร์"),
    (TEAL,      "4", "กรอกรายละเอียดปัญหา"),
    (GREEN,     "5", "กด แจ้งซ่อม → สร้าง Record"),
    (ORANGE,    "6", "ช่างกรอก สาเหตุ / วิธีแก้ไข"),
    (GREEN,     "7", "เปลี่ยนสถานะ → เสร็จแล้ว"),
]
for i, (col, num, txt) in enumerate(steps):
    y = 1.35 + i * 0.69
    add_rect(sl, 2.55, y, 0.38, 0.38, fill=col)
    add_text(sl, num, 2.55, y + 0.04, 0.38, 0.3,
             size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, txt, 3.02, y + 0.05, 2.8, 0.3, size=10, color=GRAY_DARK)
    if i < len(steps) - 1:
        add_rect(sl, 2.7, y + 0.38, 0.08, 0.31, fill=GRAY_MID)

# right: table mock + status
add_rect(sl, 6.15, 0.75, 6.98, 0.4, fill=WHITE,
         line=RGBColor(0xE5, 0xE7, 0xEB), line_w=0.5)

# filter chips
chip_data = [
    ("ทั้งหมด", True), ("รอซ่อม", False),
    ("กำลังซ่อม", False), ("เสร็จแล้ว", False),
]
cx = 6.15
for lbl, active in chip_data:
    cw = len(lbl) * 0.17 + 0.5
    add_rect(sl, cx, 1.2, cw, 0.36,
             fill=BLUE_MID if active else WHITE,
             line=BLUE_MID if active else RGBColor(0xD1, 0xD5, 0xDB),
             line_w=0.8)
    add_text(sl, lbl, cx + 0.08, 1.26, cw - 0.1, 0.25,
             size=9, bold=active, color=WHITE if active else GRAY_MID,
             align=PP_ALIGN.CENTER)
    cx += cw + 0.15

# table header
add_rect(sl, 6.15, 1.68, 6.98, 0.36, fill=GRAY_LIGHT)
cols_rl = [
    ("เครื่องแอร์", 6.25, 1.5), ("สถานที่", 7.85, 1.7),
    ("ปัญหา", 9.6, 1.8), ("สถานะ", 11.45, 0.9), ("ค่าใช้จ่าย", 12.4, 0.6),
]
for hdr, hx, hw in cols_rl:
    add_text(sl, hdr, hx, 1.72, hw, 0.27,
             size=8, bold=True, color=GRAY_MID)

rl_rows = [
    ("AC-001", "แผนกอายุรกรรม", "เสียงดัง ลมไม่เย็น", "รอซ่อม",     RED,    "-"),
    ("AC-045", "ห้องผ่าตัด",    "น้ำหยด",             "กำลังซ่อม",  ORANGE, "-"),
    ("AC-112", "ห้องฉุกเฉิน",   "แอร์ไม่ติด",         "เสร็จแล้ว",  GREEN,  "฿350"),
]
for ri, (ac, loc, prob, stat, sc, cost) in enumerate(rl_rows):
    yy = 2.1 + ri * 0.5
    add_text(sl, ac,   6.25, yy, 1.5, 0.35, size=9, bold=True, color=GRAY_DARK)
    add_text(sl, loc,  7.85, yy, 1.7, 0.35, size=9, color=GRAY_DARK)
    add_text(sl, prob, 9.6,  yy, 1.8, 0.35, size=9, color=GRAY_DARK)
    add_rect(sl, 11.45, yy + 0.04, 0.85, 0.26, fill=RGBColor(
        0xFF if sc==RED else (0xFF if sc==ORANGE else 0x22),
        0xBB if sc==RED else (0xD7 if sc==ORANGE else 0xC5),
        0xBB if sc==RED else (0x6A if sc==ORANGE else 0x5E),
    ))
    add_text(sl, stat, 11.47, yy + 0.06, 0.8, 0.22,
             size=7, bold=True, color=sc, align=PP_ALIGN.CENTER)
    add_text(sl, cost, 12.4, yy, 0.6, 0.35, size=9, color=GRAY_DARK)
    add_rect(sl, 6.15, yy + 0.42, 6.98, 0.01,
             fill=RGBColor(0xF3, 0xF4, 0xF6))

# modal mock on right side overlay
add_rect(sl, 6.15, 3.7, 6.98, 2.95, fill=RGBColor(0x00,0x00,0x00))
modal = add_rect(sl, 6.35, 3.78, 5.4, 2.72, fill=WHITE,
                 line=RGBColor(0xE5,0xE7,0xEB), line_w=1)
add_text(sl, "แก้ไขรายการซ่อม", 6.55, 3.85, 4.5, 0.38,
         size=12, bold=True, color=GRAY_DARK)
for fi, (lbl, ph) in enumerate([
    ("สาเหตุ", "ระบุสาเหตุ..."),
    ("วิธีแก้ไข", "ระบุวิธีแก้ไข..."),
    ("ค่าใช้จ่าย (บาท)", "0"),
]):
    fy = 4.3 + fi * 0.72
    add_text(sl, lbl, 6.55, fy, 2.5, 0.25, size=9, bold=True, color=GRAY_DARK)
    add_rect(sl, 6.55, fy + 0.26, 4.9, 0.38,
             fill=WHITE, line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
    add_text(sl, ph, 6.7, fy + 0.32, 4.5, 0.26, size=9, color=GRAY_MID)
add_rect(sl, 6.55, 6.1, 2.3, 0.38, fill=GRAY_LIGHT,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "ยกเลิก", 6.55, 6.16, 2.3, 0.26,
         size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
add_rect(sl, 9.05, 6.1, 2.3, 0.38, fill=BLUE_MID)
add_text(sl, "บันทึก", 9.05, 6.16, 2.3, 0.26,
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PM Schedule (overview)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("PM Schedule", "แผนซ่อมบำรุงตามรอบ — ป้องกันปัญหาก่อนเกิด")

# 4 stat cards
pm_cards = [
    ("เลยกำหนด",    "8",  RED,      "⚠"),
    ("ใกล้ถึงกำหนด","15", ORANGE,   "⏰"),
    ("ปกติ",         "210",GREEN,   "✅"),
    ("ยังไม่กำหนด",  "15", GRAY_MID,"❓"),
]
for i, (lbl, val, col, icon) in enumerate(pm_cards):
    stat_card(sl, 2.35 + i * 2.72, 0.75, 2.55, 1.2, lbl, val, col, icon)

# filter bar
add_rect(sl, 2.35, 2.1, 10.8, 0.5, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.5)
add_rect(sl, 2.4, 2.18, 2.8, 0.35, fill=WHITE,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "🔍  ค้นหา รหัส / ชื่อ...", 2.55, 2.24, 2.6, 0.25,
         size=9, color=GRAY_MID)
add_rect(sl, 5.35, 2.18, 2.0, 0.35, fill=WHITE,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "โรงพยาบาล ▾", 5.5, 2.24, 1.8, 0.25, size=9, color=GRAY_MID)
pm_chips = [("ทั้งหมด",True),("เลยกำหนด",False),("ใกล้กำหนด",False),("ปกติ",False),("ไม่มีวันที่",False)]
cx2 = 7.5
for lbl, act in pm_chips:
    cw = len(lbl) * 0.17 + 0.4
    add_rect(sl, cx2, 2.18, cw, 0.35,
             fill=BLUE_MID if act else WHITE,
             line=BLUE_MID if act else RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
    add_text(sl, lbl, cx2 + 0.06, 2.24, cw - 0.1, 0.25,
             size=9, bold=act, color=WHITE if act else GRAY_MID,
             align=PP_ALIGN.CENTER)
    cx2 += cw + 0.12

# PM table
add_rect(sl, 2.35, 2.65, 10.8, 0.4, fill=GRAY_LIGHT)
pm_headers = [("☐", 2.4, 0.3), ("รหัสแอร์", 2.75, 1.4), ("ประเภท", 4.2, 1.0),
              ("สถานที่", 5.25, 2.2), ("โรงพยาบาล", 7.5, 1.5),
              ("วัน PM ถัดไป", 9.05, 1.3), ("เหลือ", 10.4, 0.6), ("สถานะ", 11.05, 1.5)]
for hdr, hx, hw in pm_headers:
    add_text(sl, hdr, hx, 2.7, hw, 0.28, size=8, bold=True, color=GRAY_MID)

pm_rows = [
    ("☑", "AC-003", "Major","ICU / ชั้น 3 / อาคาร A", "รพ.กรุงเทพ",  "01/04/2026", "-25 วัน", "เลยกำหนด",    RED),
    ("☑", "AC-007", "Minor","แผนกศัลยกรรม / ชั้น 2","รพ.กรุงเทพ",   "15/05/2026", "-11 วัน", "เลยกำหนด",    RED),
    ("☐", "AC-045", "Fan",  "ห้องฉุกเฉิน / ชั้น 1", "รพ.สมิติเวช",  "10/06/2026", "15 วัน",  "ใกล้กำหนด",  ORANGE),
    ("☐", "AC-112", "Major","ห้องผ่าตัด / ชั้น 4",  "รพ.วิภาวดี",   "25/08/2026", "91 วัน",  "ปกติ",         GREEN),
]
for ri, (chk, ac, typ, loc, hosp, nxt, days, stat, sc) in enumerate(pm_rows):
    yy = 3.1 + ri * 0.55
    add_text(sl, chk, 2.4,  yy, 0.3, 0.38, size=10,
             color=BLUE_MID if "☑" in chk else GRAY_MID)
    add_text(sl, ac,   2.75, yy, 1.4, 0.38, size=9, bold=True, color=GRAY_DARK)
    add_text(sl, typ,  4.2,  yy, 1.0, 0.38, size=9, color=GRAY_DARK)
    add_text(sl, loc,  5.25, yy, 2.2, 0.38, size=9, color=GRAY_DARK)
    add_text(sl, hosp, 7.5,  yy, 1.5, 0.38, size=9, color=GRAY_MID)
    add_text(sl, nxt,  9.05, yy, 1.3, 0.38, size=9, color=GRAY_DARK)
    add_text(sl, days, 10.4, yy, 0.6, 0.38, size=9, bold=True, color=sc)
    add_rect(sl, 11.05, yy + 0.06, 1.0, 0.26,
             fill=RGBColor(0xFF,0xED,0xED) if sc==RED else
                  (RGBColor(0xFF,0xF7,0xED) if sc==ORANGE else
                   RGBColor(0xDC,0xFC,0xE7)))
    add_text(sl, stat, 11.08, yy + 0.09, 0.95, 0.2,
             size=7, bold=True, color=sc, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.35, yy + 0.47, 10.8, 0.01,
             fill=RGBColor(0xF3,0xF4,0xF6))

# action bar
add_rect(sl, 2.35, 5.38, 10.8, 0.52, fill=RGBColor(0xEF,0xF6,0xFF),
         line=BLUE_MID, line_w=0.8)
add_text(sl, "เลือกแล้ว 2 เครื่อง  (รพ.กรุงเทพ)",
         2.5, 5.47, 7.0, 0.32, size=10, color=BLUE_DARK, bold=True)
add_rect(sl, 10.0, 5.44, 2.7, 0.4, fill=BLUE_MID)
add_text(sl, "+ สร้างใบงาน PM", 10.0, 5.5, 2.7, 0.28,
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# note
add_text(sl, "* เลือกได้เฉพาะเครื่องในโรงพยาบาลเดียวกัน  |  กด Summary Card เพื่อ filter ทันที",
         2.35, 6.02, 10.8, 0.3, size=9, color=GRAY_MID, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — PM: สร้างใบงาน
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("PM Schedule — สร้างใบงาน", "จาก PM list → สร้าง Work Order อัตโนมัติ")

# Steps
add_text(sl, "ขั้นตอน", 2.4, 0.8, 3.0, 0.38,
         size=13, bold=True, color=BLUE_DARK)

flow_steps = [
    (BLUE_MID, "เลือกเครื่องแอร์", "Checkbox ☑ เครื่องที่ต้องการ (โรงพยาบาลเดียวกัน)"),
    (BLUE_MID, "กด + สร้างใบงาน",  "Action bar ปรากฏ → กดปุ่ม"),
    (TEAL,     "เลือกประเภทงาน",   "Major / Minor (ล้างย่อย) / Fan"),
    (TEAL,     "เลือกช่าง",        "Tech 1 (บังคับ) + Tech 2 (ถ้ามี)"),
    (GREEN,    "ยืนยัน",           "ระบบสร้าง WO + เพิ่มรายการ AC ทุกตัว → redirect ใบงาน"),
]
for i, (col, title, desc) in enumerate(flow_steps):
    y = 1.3 + i * 1.0
    add_rect(sl, 2.4, y, 0.5, 0.5, fill=col)
    add_text(sl, str(i+1), 2.4, y + 0.1, 0.5, 0.3,
             size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, title, 3.1, y + 0.05, 3.0, 0.3,
             size=11, bold=True, color=GRAY_DARK)
    add_text(sl, desc,  3.1, y + 0.32, 4.5, 0.25,
             size=9, color=GRAY_MID)
    if i < len(flow_steps)-1:
        add_rect(sl, 2.62, y + 0.5, 0.06, 0.5, fill=GRAY_MID)

# modal mock
add_rect(sl, 7.8, 0.75, 5.3, 6.55, fill=RGBColor(0x00,0x00,0x00))
add_rect(sl, 8.0, 0.85, 4.9, 6.3, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=1)
add_rect(sl, 8.0, 0.85, 4.9, 0.08, fill=BLUE_MID)
add_text(sl, "สร้างใบงาน PM", 8.2, 0.98, 4.3, 0.38,
         size=13, bold=True, color=GRAY_DARK)

# selected summary box
add_rect(sl, 8.2, 1.45, 4.5, 1.1, fill=GRAY_LIGHT,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.5)
add_text(sl, "เครื่องที่เลือก (2 เครื่อง — รพ.กรุงเทพ)",
         8.35, 1.52, 4.1, 0.3, size=9, bold=True, color=BLUE_DARK)
add_text(sl, "  •  AC-003  ICU / ชั้น 3\n  •  AC-007  แผนกศัลยกรรม / ชั้น 2",
         8.35, 1.82, 4.1, 0.6, size=9, color=GRAY_DARK)

# type select
add_text(sl, "ประเภทงาน *", 8.2, 2.66, 2.5, 0.28,
         size=9, bold=True, color=GRAY_DARK)
add_rect(sl, 8.2, 2.94, 4.5, 0.4, fill=WHITE,
         line=BLUE_MID, line_w=1)
add_text(sl, "Major (ล้างใหญ่)  ▾", 8.35, 3.01, 4.1, 0.28,
         size=10, color=GRAY_DARK)

# tech1
add_text(sl, "ช่างเทคนิค 1 *", 8.2, 3.46, 2.5, 0.28,
         size=9, bold=True, color=GRAY_DARK)
add_rect(sl, 8.2, 3.74, 4.5, 0.4, fill=WHITE,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "เลือกช่าง ▾", 8.35, 3.81, 4.1, 0.28, size=10, color=GRAY_MID)

# tech2
add_text(sl, "ช่างเทคนิค 2", 8.2, 4.26, 2.5, 0.28,
         size=9, bold=True, color=GRAY_DARK)
add_rect(sl, 8.2, 4.54, 4.5, 0.4, fill=WHITE,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "เลือกช่าง (ไม่บังคับ) ▾", 8.35, 4.61, 4.1, 0.28,
         size=10, color=GRAY_MID)

# buttons
add_rect(sl, 8.2, 5.2, 2.1, 0.42, fill=GRAY_LIGHT,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "ยกเลิก", 8.2, 5.27, 2.1, 0.28,
         size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
add_rect(sl, 10.5, 5.2, 2.2, 0.42, fill=BLUE_MID)
add_text(sl, "✓  สร้างใบงาน", 10.5, 5.27, 2.2, 0.28,
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text(sl, "→ ระบบสร้าง WO + เพิ่ม AC items + redirect ไปหน้าใบงาน",
         8.2, 5.78, 4.5, 0.28, size=8, color=GRAY_MID, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — ใบงาน: รายการ
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("ใบงาน (Work Orders) — รายการ", "ดู / กรอง / ค้นหาใบงานทั้งหมด")

# filter + search row
add_rect(sl, 2.35, 0.75, 10.8, 0.5, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.5)
add_rect(sl, 2.45, 0.82, 3.0, 0.36, fill=WHITE,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "🔍  ค้นหา เลขที่ / โรงพยาบาล...",
         2.6, 0.88, 2.8, 0.24, size=9, color=GRAY_MID)

wo_chips = [("ทั้งหมด",True),("กำลังทำ",False),("รออนุมัติ",False),("อนุมัติ",False),("ไม่อนุมัติ",False)]
cx3 = 5.6
for lbl, act in wo_chips:
    cw = len(lbl) * 0.175 + 0.35
    add_rect(sl, cx3, 0.82, cw, 0.36,
             fill=BLUE_MID if act else WHITE,
             line=BLUE_MID if act else RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
    add_text(sl, lbl, cx3 + 0.07, 0.88, cw - 0.1, 0.24,
             size=9, bold=act, color=WHITE if act else GRAY_MID,
             align=PP_ALIGN.CENTER)
    cx3 += cw + 0.12

# table
add_rect(sl, 2.35, 1.35, 10.8, 0.38, fill=GRAY_LIGHT)
wo_hdrs = [
    ("เลขที่ใบงาน", 2.45, 1.4), ("โรงพยาบาล", 4.0, 1.8),
    ("ประเภท", 5.85, 1.15), ("สถานะ", 7.05, 1.4),
    ("ช่าง", 8.5, 1.5), ("รายการ", 10.05, 0.8),
    ("วันที่", 10.9, 1.2),
]
for hdr, hx, hw in wo_hdrs:
    add_text(sl, hdr, hx, 1.38, hw, 0.3,
             size=8, bold=True, color=GRAY_MID)

wo_data = [
    ("WO-2026-001", "รพ.กรุงเทพ", "Major",   "กำลังทำ",    BLUE_MID,  "tech1",  "3", "25/05/26"),
    ("WO-2026-002", "รพ.สมิติเวช", "Minor",   "รออนุมัติ",  YELLOW,    "tech2",  "1", "24/05/26"),
    ("WO-2026-003", "รพ.วิภาวดี",  "Fan",     "อนุมัติ",    GREEN,     "tech1",  "5", "22/05/26"),
    ("WO-2026-004", "รพ.กรุงเทพ", "Major",   "ไม่อนุมัติ", RED,       "tech1",  "2", "20/05/26"),
]
for ri, (wno, hosp, typ, stat, sc, tech, items, date) in enumerate(wo_data):
    yy = 1.8 + ri * 0.65
    add_text(sl, wno,  2.45, yy, 1.5, 0.42, size=9, bold=True, color=BLUE_MID)
    add_text(sl, hosp, 4.0,  yy, 1.8, 0.42, size=9, color=GRAY_DARK)
    # type badge
    tc = BLUE_MID if typ=="Major" else (TEAL if typ=="Minor" else GREEN)
    add_rect(sl, 5.85, yy+0.06, 1.0, 0.28,
             fill=RGBColor(0xDB,0xEA,0xFE) if tc==BLUE_MID else
                  (RGBColor(0xCC,0xFB,0xF1) if tc==TEAL else
                   RGBColor(0xDC,0xFC,0xE7)))
    add_text(sl, typ, 5.88, yy+0.09, 0.95, 0.22,
             size=8, bold=True, color=tc, align=PP_ALIGN.CENTER)
    # status badge
    add_rect(sl, 7.05, yy+0.06, 1.25, 0.28,
             fill=RGBColor(0xFF,0xF7,0xED) if sc==YELLOW else
                  (RGBColor(0xDC,0xFC,0xE7) if sc==GREEN else
                   (RGBColor(0xFF,0xED,0xED) if sc==RED else
                    RGBColor(0xDB,0xEA,0xFE))))
    add_text(sl, stat, 7.08, yy+0.09, 1.2, 0.22,
             size=8, bold=True, color=sc, align=PP_ALIGN.CENTER)
    add_text(sl, tech,  8.5,  yy, 1.5, 0.42, size=9, color=GRAY_DARK)
    add_text(sl, items, 10.05,yy, 0.8, 0.42, size=9, color=GRAY_MID,
             align=PP_ALIGN.CENTER)
    add_text(sl, date,  10.9, yy, 1.2, 0.42, size=9, color=GRAY_MID)
    add_rect(sl, 2.35, yy+0.55, 10.8, 0.01,
             fill=RGBColor(0xF3,0xF4,0xF6))

# note
add_text(sl, "• คลิกแถวเพื่อเปิดรายละเอียด  •  กด + สร้างใบงาน ที่ header เพื่อเพิ่มใหม่",
         2.35, 4.6, 10.8, 0.35, size=10, color=GRAY_MID, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — ใบงาน: รายละเอียด + AC Items
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("ใบงาน — รายละเอียด (Work Order Detail)", "บันทึกงาน / ถ่ายรูป / วัดค่า / อนุมัติ")

# left panel — WO info
add_rect(sl, 2.35, 0.75, 3.8, 6.5, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.8)
add_rect(sl, 2.35, 0.75, 3.8, 0.08, fill=BLUE_MID)
add_text(sl, "ข้อมูลใบงาน", 2.5, 0.88, 3.4, 0.35,
         size=11, bold=True, color=GRAY_DARK)
info_items = [
    ("เลขที่",    "WO-2026-001"),
    ("โรงพยาบาล", "รพ.กรุงเทพ"),
    ("ประเภท",    "Major (ล้างใหญ่)"),
    ("สถานะ",     "กำลังดำเนินการ"),
    ("ช่าง 1",   "tech1"),
    ("ช่าง 2",   "tech2"),
    ("วันที่",    "25/05/2026"),
    ("หมายเหตุ",  "PM ประจำไตรมาส"),
]
for i, (k, v) in enumerate(info_items):
    y = 1.32 + i * 0.6
    add_text(sl, k, 2.5, y, 1.3, 0.3, size=9, color=GRAY_MID)
    add_text(sl, v, 3.85, y, 2.1, 0.3, size=9, bold=True, color=GRAY_DARK)

# signature area
add_rect(sl, 2.5, 6.0, 1.4, 0.9, fill=GRAY_LIGHT,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.5)
add_text(sl, "ลายเซ็น\nช่าง 1", 2.5, 6.3, 1.4, 0.4,
         size=8, color=GRAY_MID, align=PP_ALIGN.CENTER)
add_rect(sl, 4.05, 6.0, 1.4, 0.9, fill=GRAY_LIGHT,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.5)
add_text(sl, "ลายเซ็น\nผู้ตรวจ", 4.05, 6.3, 1.4, 0.4,
         size=8, color=GRAY_MID, align=PP_ALIGN.CENTER)

# right panel — AC items list
add_rect(sl, 6.35, 0.75, 7.0, 6.5, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.8)
add_text(sl, "รายการเครื่องแอร์ (3 เครื่อง)", 6.5, 0.88, 5.0, 0.35,
         size=11, bold=True, color=GRAY_DARK)

# add ac btn
add_rect(sl, 10.5, 0.84, 2.6, 0.36, fill=GRAY_LIGHT,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "+ เพิ่มเครื่องแอร์", 10.5, 0.9, 2.6, 0.24,
         size=9, color=BLUE_MID, align=PP_ALIGN.CENTER)

ac_items = [
    ("AC-001", "ห้องผ่าตัด 1 / ชั้น 4", "5 รูป", "✓ ครบ"),
    ("AC-002", "ห้องผ่าตัด 2 / ชั้น 4", "2 รูป", "○ บางส่วน"),
    ("AC-003", "ICU / ชั้น 3",           "0 รูป", "○ ยังไม่ทำ"),
]
for ri, (ac, loc, photos, status) in enumerate(ac_items):
    yy = 1.38 + ri * 1.55
    add_rect(sl, 6.45, yy, 6.75, 1.4, fill=GRAY_LIGHT,
             line=RGBColor(0xE5,0xE7,0xEB), line_w=0.5)
    add_text(sl, ac,  6.62, yy + 0.12, 2.0, 0.33, size=11, bold=True, color=GRAY_DARK)
    add_text(sl, loc, 6.62, yy + 0.45, 3.5, 0.28, size=9, color=GRAY_MID)
    add_text(sl, "📷 " + photos, 6.62, yy + 0.78, 1.5, 0.28, size=9, color=GRAY_MID)
    sc2 = GREEN if "ครบ" in status else ORANGE
    add_rect(sl, 9.5, yy + 0.68, 2.0, 0.3,
             fill=RGBColor(0xDC,0xFC,0xE7) if sc2==GREEN else RGBColor(0xFF,0xF7,0xED))
    add_text(sl, status, 9.55, yy + 0.74, 1.9, 0.2,
             size=8, bold=True, color=sc2, align=PP_ALIGN.CENTER)
    add_rect(sl, 11.6, yy + 0.68, 1.3, 0.3, fill=BLUE_MID)
    add_text(sl, "เปิดดูรายละเอียด", 11.62, yy + 0.74, 1.25, 0.2,
             size=7, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# approve buttons at bottom
add_rect(sl, 6.45, 6.0, 3.1, 0.5, fill=GRAY_LIGHT,
         line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
add_text(sl, "ส่งอนุมัติ (Tech)", 6.45, 6.1, 3.1, 0.3,
         size=10, color=GRAY_DARK, align=PP_ALIGN.CENTER)
add_rect(sl, 9.75, 6.0, 1.8, 0.5, fill=GREEN)
add_text(sl, "✓ อนุมัติ", 9.75, 6.1, 1.8, 0.3,
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_rect(sl, 11.75, 6.0, 1.4, 0.5, fill=RED)
add_text(sl, "✗ ไม่อนุมัติ", 11.75, 6.1, 1.4, 0.3,
         size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "* ปุ่มอนุมัติ / ไม่อนุมัติ แสดงเฉพาะ Owner + Admin",
         6.45, 6.6, 6.7, 0.28, size=8, color=GRAY_MID, italic=True)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — AC Item Detail (Photos / Checklist / Measurements)
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("รายละเอียดเครื่องแอร์ (AC Item Detail)", "ถ่ายรูป / บันทึกค่าวัด / Checklist")

# tabs
tabs = [("รูปภาพ", True), ("ค่าวัด (Major)", False), ("Checklist", False)]
tx = 2.4
for lbl, act in tabs:
    tw = 2.2
    add_rect(sl, tx, 0.75, tw, 0.45,
             fill=BLUE_MID if act else GRAY_LIGHT,
             line=BLUE_MID, line_w=0.8)
    add_text(sl, lbl, tx, 0.84, tw, 0.28,
             size=10, bold=act, color=WHITE if act else GRAY_MID,
             align=PP_ALIGN.CENTER)
    tx += tw + 0.05

# photos tab content
add_text(sl, "Tab: รูปภาพ", 2.4, 1.3, 5.0, 0.35,
         size=11, bold=True, color=BLUE_DARK)
# photo grid mock
for pi in range(6):
    px = 2.4 + (pi % 3) * 1.85
    py = 1.75 + (pi // 3) * 1.65
    add_rect(sl, px, py, 1.65, 1.4, fill=GRAY_LIGHT,
             line=RGBColor(0xE5,0xE7,0xEB), line_w=0.5)
    if pi < 4:
        add_text(sl, "📷", px + 0.6, py + 0.45, 0.6, 0.5, size=20, color=GRAY_MID)
    else:
        add_rect(sl, px, py, 1.65, 1.4, fill=RGBColor(0xDB,0xEA,0xFE),
                 line=BLUE_MID, line_w=0.8)
        add_text(sl, "📷 +\nถ่ายรูป", px + 0.4, py + 0.45, 0.9, 0.55,
                 size=10, color=BLUE_MID, align=PP_ALIGN.CENTER)

# measurements panel
add_rect(sl, 8.0, 0.75, 5.2, 6.5, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.8)
add_text(sl, "Tab: ค่าวัด (Major only)", 8.2, 0.85, 4.7, 0.35,
         size=11, bold=True, color=BLUE_DARK)

meas_fields = [
    ("อุณหภูมิลมออก (°C)", "°C"),
    ("อุณหภูมิลมกลับ (°C)", "°C"),
    ("กระแสไฟ Comp (A)", "A"),
    ("กระแสไฟ Fan (A)", "A"),
    ("แรงดันไฟ (V)", "V"),
    ("แรงดัน Low Side (PSI)", "PSI"),
    ("แรงดัน High Side (PSI)", "PSI"),
]
for mi, (fname, unit) in enumerate(meas_fields):
    my = 1.35 + mi * 0.68
    add_text(sl, fname, 8.2, my, 3.2, 0.28, size=9, color=GRAY_DARK)
    add_rect(sl, 11.45, my, 1.4, 0.38, fill=WHITE,
             line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
    add_text(sl, unit, 11.6, my + 0.07, 1.1, 0.24,
             size=9, color=GRAY_MID, align=PP_ALIGN.RIGHT)

add_text(sl, "Tab: Checklist", 8.2, 6.15, 4.0, 0.3,
         size=10, bold=True, color=GRAY_MID)
checklist_items = ["ทำความสะอาด Filter", "ล้างคอยล์เย็น", "ตรวจสอบสายไฟ"]
for ci, item in enumerate(checklist_items):
    add_text(sl, "☑  " + item, 8.35, 6.48 + ci * 0.01, 4.5, 0.28,
             size=9, color=GRAY_MID)
add_text(sl, "☑  ทำความสะอาด Filter   ☑  ล้างคอยล์เย็น   ☑  ตรวจสอบสายไฟ",
         8.2, 6.5, 4.7, 0.28, size=8, color=GRAY_MID)

# save button
add_rect(sl, 8.2, 6.88, 4.7, 0.42, fill=BLUE_MID)
add_text(sl, "บันทึกค่า", 8.2, 6.94, 4.7, 0.28,
         size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Master Data
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("Master Data", "จัดการข้อมูลหลัก — เฉพาะ Admin + Owner")

# tabs
tab_labels = [("Hospitals", True), ("อาคาร / ชั้น / แผนก", False), ("เครื่องแอร์", False)]
tx2 = 2.4
for lbl, act in tab_labels:
    tw = 3.4
    add_rect(sl, tx2, 0.75, tw, 0.42,
             fill=BLUE_MID if act else GRAY_LIGHT,
             line=BLUE_MID, line_w=0.8)
    add_text(sl, lbl, tx2, 0.84, tw, 0.26,
             size=10, bold=act, color=WHITE if act else GRAY_MID,
             align=PP_ALIGN.CENTER)
    tx2 += tw + 0.05

# hospitals tab mock
add_rect(sl, 2.35, 1.27, 10.8, 0.38, fill=GRAY_LIGHT)
for hdr, hx, hw in [("ชื่อโรงพยาบาล", 2.5, 3.0), ("รหัส/โน้ต", 5.55, 2.5),
                     ("สถานะ", 8.1, 1.2), ("", 9.35, 2.5)]:
    add_text(sl, hdr, hx, 1.3, hw, 0.3, size=8, bold=True, color=GRAY_MID)

hosp_rows = [
    ("โรงพยาบาลกรุงเทพ",  "HBK-01", "ใช้งาน", GREEN),
    ("โรงพยาบาลสมิติเวช",  "SMT-01", "ใช้งาน", GREEN),
    ("โรงพยาบาลวิภาวดี",   "VPW-01", "ใช้งาน", GREEN),
    ("โรงพยาบาลบำรุงราษฎร์","BRR-01", "ปิดใช้",  RED),
]
for ri, (name, code, stat, sc) in enumerate(hosp_rows):
    yy = 1.72 + ri * 0.55
    add_text(sl, name, 2.5,  yy, 3.0, 0.38, size=10, color=GRAY_DARK)
    add_text(sl, code, 5.55, yy, 2.5, 0.38, size=9, color=GRAY_MID)
    add_rect(sl, 8.1, yy+0.06, 1.0, 0.26,
             fill=RGBColor(0xDC,0xFC,0xE7) if sc==GREEN else RGBColor(0xFF,0xED,0xED))
    add_text(sl, stat, 8.13, yy+0.09, 0.95, 0.2,
             size=8, bold=True, color=sc, align=PP_ALIGN.CENTER)
    add_rect(sl, 9.35, yy+0.04, 0.8, 0.3, fill=GRAY_LIGHT,
             line=RGBColor(0xD1,0xD5,0xDB), line_w=0.5)
    add_text(sl, "✏ แก้ไข", 9.38, yy+0.08, 0.75, 0.22,
             size=8, color=BLUE_MID, align=PP_ALIGN.CENTER)
    add_rect(sl, 10.3, yy+0.04, 0.8, 0.3, fill=GRAY_LIGHT,
             line=RGBColor(0xD1,0xD5,0xDB), line_w=0.5)
    add_text(sl, "🗑 ลบ", 10.33, yy+0.08, 0.75, 0.22,
             size=8, color=RED, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.35, yy+0.47, 10.8, 0.01,
             fill=RGBColor(0xF3,0xF4,0xF6))

# structure tab description
add_rect(sl, 2.35, 4.1, 10.8, 2.55, fill=GRAY_LIGHT,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.5)
add_text(sl, "Tab: อาคาร / ชั้น / แผนก  —  Drill-down 3 คอลัมน์",
         2.5, 4.2, 8.0, 0.35, size=11, bold=True, color=BLUE_DARK)

# mock 3-col drill down
col_titles = ["อาคาร", "ชั้น", "แผนก"]
col_items  = [
    ["อาคาร A", "อาคาร B"],
    ["ชั้น 1", "ชั้น 2", "ชั้น 3"],
    ["แผนกอายุรกรรม", "แผนกศัลยกรรม", "ICU"],
]
for ci, (ct, items) in enumerate(zip(col_titles, col_items)):
    cx4 = 2.45 + ci * 3.55
    add_rect(sl, cx4, 4.65, 3.3, 0.32, fill=RGBColor(0xDB,0xEA,0xFE))
    add_text(sl, ct, cx4 + 0.1, 4.7, 3.1, 0.22,
             size=9, bold=True, color=BLUE_DARK)
    for ii, item in enumerate(items):
        iy = 5.05 + ii * 0.36
        bg = RGBColor(0xDB,0xEA,0xFE) if (ci == 0 and ii == 0) else WHITE
        add_rect(sl, cx4, iy, 3.3, 0.3,
                 fill=bg, line=RGBColor(0xE5,0xE7,0xEB), line_w=0.3)
        add_text(sl, item, cx4 + 0.15, iy + 0.05, 2.8, 0.22,
                 size=9, color=BLUE_DARK if bg!=WHITE else GRAY_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ผู้ใช้งาน
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("ผู้ใช้งาน (Users)", "จัดการบัญชีผู้ใช้ — เฉพาะ Admin")

# table
add_rect(sl, 2.35, 0.75, 10.8, 0.4, fill=GRAY_LIGHT)
u_hdrs = [("", 2.4, 0.5), ("ชื่อ", 2.95, 2.0), ("Username", 5.0, 1.7),
          ("Role", 6.75, 1.3), ("สถานะ", 8.1, 1.2), ("", 9.35, 2.0)]
for hdr, hx, hw in u_hdrs:
    add_text(sl, hdr, hx, 0.78, hw, 0.3, size=8, bold=True, color=GRAY_MID)
add_text(sl, "ชื่อ", 2.95, 0.78, 2.0, 0.3, size=8, bold=True, color=GRAY_MID)
add_text(sl, "Username", 5.0, 0.78, 1.7, 0.3, size=8, bold=True, color=GRAY_MID)
add_text(sl, "Role", 6.75, 0.78, 1.3, 0.3, size=8, bold=True, color=GRAY_MID)
add_text(sl, "สถานะ", 8.1, 0.78, 1.2, 0.3, size=8, bold=True, color=GRAY_MID)

users_data = [
    ("A", "Administrator",  "admin",  "Admin",  BLUE_DARK, "ใช้งาน", GREEN),
    ("O", "สมชาย จันทร์ดี", "owner1", "Owner",  TEAL,      "ใช้งาน", GREEN),
    ("T", "สุดา ใจดี",      "tech1",  "Tech",   GREEN,     "ใช้งาน", GREEN),
    ("T", "มานะ รักงาน",   "tech2",  "Tech",   GREEN,     "หยุดใช้", RED),
]
for ri, (init, name, uname, role, rc, stat, sc) in enumerate(users_data):
    yy = 1.28 + ri * 0.72
    add_rect(sl, 2.4, yy+0.05, 0.42, 0.42, fill=rc)
    add_text(sl, init, 2.4, yy+0.12, 0.42, 0.28,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, name,  2.95, yy+0.1, 2.0, 0.28, size=10, color=GRAY_DARK)
    add_text(sl, uname, 5.0,  yy+0.1, 1.7, 0.28, size=9,
             color=GRAY_MID, italic=True)
    add_rect(sl, 6.75, yy+0.12, 1.1, 0.3, fill=rc)
    add_text(sl, role, 6.78, yy+0.16, 1.05, 0.22,
             size=8, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(sl, 8.1, yy+0.12, 1.0, 0.3,
             fill=RGBColor(0xDC,0xFC,0xE7) if sc==GREEN else RGBColor(0xFF,0xED,0xED))
    add_text(sl, stat, 8.13, yy+0.16, 0.95, 0.22,
             size=8, bold=True, color=sc, align=PP_ALIGN.CENTER)
    add_rect(sl, 9.35, yy+0.1, 0.8, 0.32, fill=GRAY_LIGHT,
             line=RGBColor(0xD1,0xD5,0xDB), line_w=0.5)
    add_text(sl, "✏ แก้ไข", 9.38, yy+0.15, 0.75, 0.22,
             size=8, color=BLUE_MID, align=PP_ALIGN.CENTER)
    add_rect(sl, 2.35, yy+0.64, 10.8, 0.01,
             fill=RGBColor(0xF3,0xF4,0xF6))

# create user modal mock
add_rect(sl, 6.0, 3.7, 7.1, 3.55, fill=RGBColor(0x00,0x00,0x00))
add_rect(sl, 6.2, 3.82, 6.7, 3.3, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=1)
add_rect(sl, 6.2, 3.82, 6.7, 0.08, fill=BLUE_MID)
add_text(sl, "เพิ่มผู้ใช้งานใหม่", 6.4, 3.95, 5.5, 0.35,
         size=12, bold=True, color=GRAY_DARK)
for fi, (lbl, ph) in enumerate([
    ("ชื่อ-นามสกุล *", "ชื่อเต็ม"),
    ("Username *", "username"),
    ("Password *", "••••••••"),
    ("Role *", "Admin / Owner / Tech"),
]):
    fy = 4.4 + fi * 0.6
    add_text(sl, lbl, 6.4, fy, 2.5, 0.25, size=9, bold=True, color=GRAY_DARK)
    add_rect(sl, 6.4, fy + 0.26, 6.3, 0.38,
             fill=WHITE, line=RGBColor(0xD1,0xD5,0xDB), line_w=0.8)
    add_text(sl, ph, 6.55, fy + 0.32, 5.9, 0.26, size=9, color=GRAY_MID)
add_rect(sl, 9.75, 6.7, 2.75, 0.4, fill=BLUE_MID)
add_text(sl, "✓ บันทึก", 9.75, 6.76, 2.75, 0.28,
         size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Workflow Overview
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("Workflow ภาพรวม", "กระบวนการทำงานตั้งแต่ต้นจนจบ")

flow_items = [
    (BLUE_MID, "แจ้งซ่อม\nหรือ PM Schedule",    "Tech / ทุก Role\nเลือกเครื่อง → ระบุปัญหา"),
    (TEAL,     "สร้างใบงาน\n(Work Order)",       "Tech / Admin\nเลือกประเภท + ช่าง"),
    (ORANGE,   "บันทึกงาน\n(AC Item Detail)",    "Tech\nรูปภาพ + ค่าวัด + Checklist"),
    (BLUE_MID, "ส่งอนุมัติ",                     "Tech\nกดปุ่ม Submit"),
    (GREEN,    "อนุมัติ / ไม่อนุมัติ",           "Owner / Admin\nPM Date อัปเดตอัตโนมัติ"),
    (TEAL,     "ปิดงาน",                         "ระบบ update สถานะ\nพร้อม PM รอบถัดไป"),
]

box_w = 1.7
box_h = 2.0
start_x = 2.5

for i, (col, title, desc) in enumerate(flow_items):
    bx = start_x + i * (box_w + 0.42)
    by = 1.8

    add_rect(sl, bx, by, box_w, box_h, fill=col)
    add_text(sl, title, bx + 0.1, by + 0.2, box_w - 0.2, 0.7,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, desc, bx + 0.1, by + 1.0, box_w - 0.2, 0.85,
             size=8, color=RGBColor(0xDB,0xEA,0xFE) if col==BLUE_MID else
                          (RGBColor(0xCC,0xFB,0xF1) if col==TEAL else
                           (RGBColor(0xFF,0xF7,0xED) if col==ORANGE else
                            RGBColor(0xDC,0xFC,0xE7))),
             align=PP_ALIGN.CENTER)
    # step number circle
    add_rect(sl, bx + box_w/2 - 0.22, by - 0.52, 0.44, 0.44, fill=col)
    add_text(sl, str(i+1), bx + box_w/2 - 0.22, by - 0.45, 0.44, 0.28,
             size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    if i < len(flow_items) - 1:
        ax = bx + box_w + 0.05
        add_rect(sl, ax, by + box_h/2 - 0.06, 0.32, 0.12, fill=GRAY_MID)
        add_text(sl, "▶", ax + 0.02, by + box_h/2 - 0.18, 0.3, 0.3,
                 size=10, color=GRAY_MID)

# Role responsibilities
add_rect(sl, 2.35, 4.45, 10.8, 2.85, fill=GRAY_LIGHT,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.5)
add_text(sl, "สรุปสิทธิ์ตาม Role", 2.5, 4.55, 5.0, 0.35,
         size=11, bold=True, color=BLUE_DARK)

role_data = [
    ("Admin",  BLUE_DARK, ["ทุกสิทธิ์ในระบบ", "จัดการ Users + Master Data",
                            "อนุมัติใบงาน", "ดู Dashboard ทั้งหมด"]),
    ("Owner",  TEAL,      ["ดู Dashboard", "ดู + อนุมัติใบงาน",
                            "จัดการ Master Data", "ไม่สามารถจัดการ Users"]),
    ("Tech",   GREEN,     ["แจ้งซ่อม", "สร้างใบงาน",
                            "บันทึกงาน (ถ่ายรูป/วัดค่า/Checklist)", "ส่งอนุมัติ"]),
]
for ri, (role, col, perms) in enumerate(role_data):
    rx = 2.5 + ri * 3.6
    add_rect(sl, rx, 5.0, 0.9, 0.32, fill=col)
    add_text(sl, role, rx, 5.06, 0.9, 0.22,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for pi, perm in enumerate(perms):
        add_text(sl, "•  " + perm, rx, 5.42 + pi * 0.38, 3.4, 0.32,
                 size=9, color=GRAY_DARK)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — PDF + Tips
# ══════════════════════════════════════════════════════════════════════════════
sl = new_slide("PDF Download & Tips", "ดาวน์โหลดรายงาน + เคล็ดลับการใช้งาน")

# PDF section
add_rect(sl, 2.35, 0.75, 5.0, 5.5, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.8)
add_rect(sl, 2.35, 0.75, 5.0, 0.08, fill=RED)
add_text(sl, "📄  ดาวน์โหลด PDF ใบงาน", 2.55, 0.9, 4.6, 0.38,
         size=12, bold=True, color=GRAY_DARK)
pdf_steps = [
    "เปิดหน้า Work Order Detail",
    "กดปุ่ม  📄 PDF  (มุมบนขวา)",
    "Browser เปิด Tab ใหม่",
    "PDF แสดง: ข้อมูลใบงาน + รายการ AC\n  + Checklist + ลายเซ็น",
    "กด Print หรือ Save as PDF",
]
for si, step in enumerate(pdf_steps):
    sy = 1.42 + si * 0.78
    add_rect(sl, 2.55, sy, 0.36, 0.36, fill=RED)
    add_text(sl, str(si+1), 2.55, sy+0.05, 0.36, 0.26,
             size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(sl, step, 3.0, sy+0.04, 4.1, 0.55, size=10, color=GRAY_DARK)

# Tips section
add_rect(sl, 7.55, 0.75, 5.6, 5.5, fill=WHITE,
         line=RGBColor(0xE5,0xE7,0xEB), line_w=0.8)
add_rect(sl, 7.55, 0.75, 5.6, 0.08, fill=YELLOW)
add_text(sl, "💡  เคล็ดลับการใช้งาน", 7.75, 0.9, 5.1, 0.38,
         size=12, bold=True, color=GRAY_DARK)

tips = [
    ("PM auto-update",  "เมื่ออนุมัติใบงาน → next_pm_date\nอัปเดตอัตโนมัติตาม interval"),
    ("Filter PM",       "กด Summary Card (overdue/due_soon)\nเพื่อ filter ทันที"),
    ("ค้นหาเครื่องแอร์", "แจ้งซ่อม/PM: พิมพ์ รหัส/ชื่อ/แผนก\nในช่อง search"),
    ("Mobile",          "รองรับมือถือ: Hamburger menu\nเปิด sidebar"),
    ("Multi-select PM", "เลือกหลายเครื่องแล้วสร้าง WO เดียว\n(โรงพยาบาลเดียวกันเท่านั้น)"),
    ("Session persist", "ปิดเบราว์เซอร์แล้วเปิดใหม่\nไม่ต้อง Login ซ้ำ"),
]
for ti, (t, d) in enumerate(tips):
    ty = 1.38 + ti * 0.72
    add_rect(sl, 7.75, ty, 0.08, 0.38, fill=YELLOW)
    add_text(sl, t, 7.95, ty, 2.2, 0.28, size=10, bold=True, color=GRAY_DARK)
    add_text(sl, d, 7.95, ty + 0.3, 5.0, 0.38, size=9, color=GRAY_MID)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — End
# ══════════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(BLANK)
add_rect(sl, 0, 0, 13.33, 7.5, fill=BLUE_DARK)
add_rect(sl, 0, 6.3, 13.33, 0.05, fill=BLUE_MID)
add_rect(sl, 0, 6.35, 13.33, 1.15, fill=RGBColor(0x17, 0x30, 0x52))
add_text(sl, "❄", 6.2, 1.8, 1.0, 1.0, size=52, color=RGBColor(0x60,0xA5,0xFA))
add_text(sl, "Air System", 2.0, 2.8, 9.33, 1.0,
         size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(sl, "ขอบคุณที่ใช้งานระบบ", 2.0, 3.8, 9.33, 0.6,
         size=18, color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)
add_text(sl, "Technical Water Co.,Ltd  |  Air Conditioning Maintenance System",
         2.0, 6.5, 9.33, 0.35, size=10,
         color=RGBColor(0x93,0xC5,0xFD), align=PP_ALIGN.CENTER)

# ─── SAVE ─────────────────────────────────────────────────────────────────────
out = r"C:\Users\Macha\Desktop\AirSystem_Manual.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
